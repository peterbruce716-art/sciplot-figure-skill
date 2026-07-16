from __future__ import annotations

import argparse
import json
import re
from pathlib import Path
from typing import Any

from object_reconstruction import load_json, write_json


def _vba_text(value: str) -> str:
    return '"' + str(value).replace('"', '""').replace("\r", " ").replace("\n", " ") + '"'


def _rgb(value: Any, default: str = "#FFFFFF") -> str:
    text = str(value or default).lstrip("#")
    if len(text) != 6 or not re.fullmatch(r"[0-9A-Fa-f]{6}", text):
        text = default.lstrip("#")
    r, g, b = int(text[0:2], 16), int(text[2:4], 16), int(text[4:6], 16)
    return f"RGB({r}, {g}, {b})"


def _coords(element: dict[str, Any], width: int, height: int, slide_width: float = 960.0, slide_height: float = 540.0) -> tuple[float, float, float, float]:
    x, y, w, h = [float(value) for value in element["bbox_px"]]
    scale = min(slide_width / width, slide_height / height)
    offset_x = (slide_width - width * scale) / 2
    offset_y = (slide_height - height * scale) / 2
    return offset_x + x * scale, offset_y + y * scale, w * scale, h * scale


def generate_vba(manifest: dict[str, Any], *, asset_dir_name: str = "assets") -> str:
    width = int(manifest["source"]["width_px"])
    height = int(manifest["source"]["height_px"])
    lines = [
        "Attribute VB_Name = \"SciPlotObjectReconstruction\"",
        "Option Explicit",
        "",
        "' Generated from a SciPlot Object Manifest. Inspect before enabling macros.",
        "Private Const GENERATED_PREFIX As String = \"SCIPLOT_OBJ_\"",
        "Private Const CANVAS_W As Double = 960#",
        "Private Const CANVAS_H As Double = 540#",
        "Private Const SOURCE_W As Double = %d#" % width,
        "Private Const SOURCE_H As Double = %d#" % height,
        "",
        "Private Function XPos(ByVal px As Double) As Double: XPos = (CANVAS_W - SOURCE_W * ScaleFactor()) / 2# + px * ScaleFactor(): End Function",
        "Private Function YPos(ByVal px As Double) As Double: YPos = (CANVAS_H - SOURCE_H * ScaleFactor()) / 2# + px * ScaleFactor(): End Function",
        "Private Function ScaleFactor() As Double: ScaleFactor = IIf(CANVAS_W / SOURCE_W < CANVAS_H / SOURCE_H, CANVAS_W / SOURCE_W, CANVAS_H / SOURCE_H): End Function",
        "",
        "Public Sub ClearSciPlotObjects()",
        "    Dim i As Long",
        "    For i = ActivePresentation.Slides(1).Shapes.Count To 1 Step -1",
        "        If Left$(ActivePresentation.Slides(1).Shapes(i).Name, Len(GENERATED_PREFIX)) = GENERATED_PREFIX Then ActivePresentation.Slides(1).Shapes(i).Delete",
        "    Next i",
        "End Sub",
        "",
        "Public Sub BuildSkeleton()",
        "    Dim s As Slide: Set s = ActivePresentation.Slides(1)",
    ]
    for element in sorted(manifest.get("elements", []), key=lambda item: int(item.get("z_order", 0))):
        x, y, w, h = _coords(element, width, height)
        eid = re.sub(r"[^A-Za-z0-9_]", "_", str(element.get("id")))
        primitive = element.get("primitive")
        name = f"{eid}_skeleton"
        if primitive in {"line", "arrow", "connector", "polyline"}:
            points = element.get("observed_endpoints_px") or [[x, y + h / 2], [x + w, y + h / 2]]
            p1, p2 = points[0], points[-1]
            lines.append(f"    Set shp = s.Shapes.AddConnector(msoConnectorStraight, XPos({float(p1[0])}), YPos({float(p1[1])}), XPos({float(p2[0])}), YPos({float(p2[1])}))")
        else:
            lines.append(f"    Set shp = s.Shapes.AddShape(msoShapeRectangle, XPos({x}), YPos({y}), {w}, {h})")
        lines.append(f"    shp.Name = GENERATED_PREFIX & {_vba_text(name)}: shp.Fill.ForeColor.RGB = RGB(200, 200, 200): shp.Line.ForeColor.RGB = RGB(80, 80, 80)")
        lines.append(f"    shp.TextFrame.TextRange.Text = {_vba_text(str(element.get('id')))}")
    lines.extend(["End Sub", "", "Public Sub BuildFinal()", "    Dim s As Slide: Set s = ActivePresentation.Slides(1)", "    Dim shp As Shape"])
    for element in sorted(manifest.get("elements", []), key=lambda item: int(item.get("z_order", 0))):
        x, y, w, h = _coords(element, width, height)
        eid = re.sub(r"[^A-Za-z0-9_]", "_", str(element.get("id")))
        primitive = element.get("primitive")
        style = element.get("style", {})
        if element.get("bucket") == "preserved_raster":
            asset = Path(str(element.get("asset_path", f"{eid}.png"))).as_posix()
            lines.append(f"    Set shp = s.Shapes.AddPicture({ _vba_text(asset) }, msoFalse, msoTrue, XPos({float(element['bbox_px'][0])}), YPos({float(element['bbox_px'][1])}), {w}, {h})")
        elif primitive in {"line", "arrow", "connector", "polyline"}:
            points = element.get("observed_endpoints_px") or [[float(element['bbox_px'][0]), float(element['bbox_px'][1])], [float(element['bbox_px'][0]) + float(element['bbox_px'][2]), float(element['bbox_px'][1]) + float(element['bbox_px'][3])]]
            p1, p2 = points[0], points[-1]
            lines.append(f"    Set shp = s.Shapes.AddConnector(msoConnectorStraight, XPos({float(p1[0])}), YPos({float(p1[1])}), XPos({float(p2[0])}), YPos({float(p2[1])}))")
            if primitive in {"arrow", "connector"}:
                lines.append("    shp.Line.EndArrowheadStyle = msoArrowheadTriangle")
        elif primitive == "ellipse" or primitive == "circle":
            lines.append(f"    Set shp = s.Shapes.AddShape(msoShapeOval, XPos({x}), YPos({y}), {w}, {h})")
        elif primitive == "rounded_rectangle":
            lines.append(f"    Set shp = s.Shapes.AddShape(msoShapeRoundedRectangle, XPos({x}), YPos({y}), {w}, {h})")
        else:
            lines.append(f"    Set shp = s.Shapes.AddShape(msoShapeRectangle, XPos({x}), YPos({y}), {w}, {h})")
        lines.extend([
            f"    shp.Name = GENERATED_PREFIX & {_vba_text(eid)}",
            f"    shp.Fill.ForeColor.RGB = {_rgb(style.get('fill'), '#FFFFFF')}",
            f"    shp.Line.ForeColor.RGB = {_rgb(style.get('stroke'), '#333333')}",
            f"    shp.Line.Weight = {float(style.get('stroke_width_pt', 1.0))}",
        ])
        text = str((element.get("text") or {}).get("content", ""))
        if not text:
            text = str(element.get("text_content", ""))
        if text:
            lines.append(f"    shp.TextFrame.TextRange.Text = {_vba_text(text)}")
    lines.extend(["End Sub", ""])
    return "\n".join(lines) + "\n"


def export(manifest: dict[str, Any], output_dir: Path, *, backend: str = "powerpoint_vba", create_pptx: bool = False) -> dict[str, Any]:
    output_dir.mkdir(parents=True, exist_ok=True)
    vba_path = output_dir / "reconstruct_figure.bas"
    vba_path.write_text(generate_vba(manifest), encoding="utf-8")
    unsupported = [str(item.get("id")) for item in manifest.get("elements", []) if item.get("primitive") == "unknown"]
    pptx_path: Path | None = None
    reason = "offline_code_generation_mode"
    if create_pptx:
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Inches  # type: ignore
            presentation = Presentation()
            presentation.slide_width = Inches(10)
            presentation.slide_height = Inches(5.625)
            presentation.slides.add_slide(presentation.slide_layouts[6])
            pptx_path = output_dir / "figure.pptx"
            presentation.save(pptx_path)
            reason = "pptx_container_generated; shape materialization requires runtime verification"
        except Exception as exc:
            reason = f"pptx_backend_unavailable:{type(exc).__name__}"
    editable = sum(1 for item in manifest.get("elements", []) if item.get("bucket") == "editable_vector")
    preserved = sum(1 for item in manifest.get("elements", []) if item.get("bucket") == "preserved_raster")
    report = {
        "schema": "scientificfigure.office_export_report.v1",
        "backend": backend,
        "pptx": None if pptx_path is None else pptx_path.name,
        "vba": vba_path.name,
        "exported_elements": editable + preserved,
        "editable_elements": editable,
        "preserved_raster_elements": preserved,
        "unsupported_elements": unsupported,
        "materialization_attempted": bool(create_pptx),
        "reason": reason,
        "status": "failed" if unsupported else "generated_not_runtime_verified",
    }
    return report


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate optional editable Office/VBA output without making Office a core dependency.")
    parser.add_argument("--manifest", required=True, type=Path)
    parser.add_argument("--output-dir", required=True, type=Path)
    parser.add_argument("--backend", default="powerpoint_vba")
    parser.add_argument("--create-pptx", action="store_true")
    parser.add_argument("--report", required=True, type=Path)
    args = parser.parse_args()
    result = export(load_json(args.manifest), args.output_dir, backend=args.backend, create_pptx=args.create_pptx)
    write_json(args.report, result)
    print(json.dumps(result, ensure_ascii=False, indent=2, sort_keys=True))
    return 0 if result["status"] != "failed" else 2


if __name__ == "__main__":
    raise SystemExit(main())
